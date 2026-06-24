#!/usr/bin/env python3
import sys
import csv
import json
import pathlib


OUTPUT_DIR = pathlib.Path(__file__).parent / "converted"
MAX_FILE_SIZE_MB = 500


def ensure_output_dir():
    OUTPUT_DIR.mkdir(exist_ok=True)


def convert_pdf(path):
    import pdfplumber
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def convert_docx(path):
    import mammoth
    with open(path, "rb") as f:
        result = mammoth.convert_to_markdown(f)
    return result.value


def convert_pptx(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    sections = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        sections.append("\n".join(lines))
    return "\n\n".join(sections)


def _rows_to_md_table(rows):
    if not rows:
        return ""
    header = rows[0]
    col_widths = [max(len(str(header[c])), max((len(str(r[c])) for r in rows[1:]), default=0)) for c in range(len(header))]
    def fmt_row(row):
        return "| " + " | ".join(str(row[c]).ljust(col_widths[c]) for c in range(len(header))) + " |"
    sep = "| " + " | ".join("-" * w for w in col_widths) + " |"
    lines = [fmt_row(header), sep] + [fmt_row(r) for r in rows[1:]]
    return "\n".join(lines)


def convert_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = [[cell.value if cell.value is not None else "" for cell in row] for row in sheet.iter_rows()]
        if rows:
            parts.append(f"## {sheet.title}\n\n{_rows_to_md_table(rows)}")
    return "\n\n".join(parts)


def convert_csv(path):
    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = list(reader)
    return _rows_to_md_table(rows)


def convert_html(path):
    from bs4 import BeautifulSoup
    with open(path, encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text(separator="\n")


def convert_txt(path):
    with open(path, encoding="utf-8") as f:
        return f.read()


def convert_rst(path):
    with open(path, encoding="utf-8") as f:
        content = f.read()
    lines = content.splitlines()
    output = []
    i = 0
    while i < len(lines):
        line = lines[i]
        if line.startswith(".. code") or line.startswith(".. sourcecode"):
            output.append("```")
            i += 1
            while i < len(lines) and (not lines[i].strip() or lines[i].startswith("   ")):
                output.append(lines[i][3:] if lines[i].startswith("   ") else lines[i])
                i += 1
            output.append("```")
        else:
            output.append(line)
            i += 1
    return "\n".join(output)


def convert_org(path):
    with open(path, encoding="utf-8") as f:
        return f.read()


def convert_ipynb(path):
    with open(path, encoding="utf-8") as f:
        nb = json.load(f)
    if "cells" not in nb:
        raise ValueError(f"{path.name} is not a valid Jupyter notebook (missing 'cells' key)")
    parts = []
    for cell in nb["cells"]:
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", []))
        if cell_type == "markdown":
            parts.append(source)
        elif cell_type == "code":
            parts.append(f"```python\n{source}\n```")
    return "\n\n".join(parts)


def convert_image(path):
    import pytesseract
    from PIL import Image
    img = Image.open(str(path))
    return pytesseract.image_to_string(img)


CONVERTERS = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".pptx": convert_pptx,
    ".xlsx": convert_xlsx,
    ".csv": convert_csv,
    ".html": convert_html,
    ".htm": convert_html,
    ".txt": convert_txt,
    ".rst": convert_rst,
    ".org": convert_org,
    ".ipynb": convert_ipynb,
    ".png": convert_image,
    ".jpg": convert_image,
    ".jpeg": convert_image,
}


def convert_file(file_path):
    path = pathlib.Path(file_path)

    if not path.exists():
        print(f"Not found: {file_path}", file=sys.stderr)
        return

    ext = path.suffix.lower()
    if ext not in CONVERTERS:
        print(f"Unsupported: {path.name}")
        return

    size_mb = path.stat().st_size / (1024 * 1024)
    if size_mb > MAX_FILE_SIZE_MB:
        print(f"Skipped {path.name}: file too large ({size_mb:.1f} MB > {MAX_FILE_SIZE_MB} MB)", file=sys.stderr)
        return

    out_name = path.stem + "_" + ext.lstrip(".") + ".md"
    out_path = OUTPUT_DIR / out_name

    if out_path.exists():
        print(f"Warning: overwriting {out_path}", file=sys.stderr)

    try:
        content = CONVERTERS[ext](path)
        out_path.write_text(content or "", encoding="utf-8")
        print(f"Converted: {path.name} -> {out_path}")
    except Exception as e:
        print(f"Error converting {path.name}: {e}", file=sys.stderr)


def main():
    if len(sys.argv) < 2:
        print("Usage: convert-to-md file1 file2 ...")
        sys.exit(1)
    ensure_output_dir()
    for arg in sys.argv[1:]:
        convert_file(arg)


if __name__ == "__main__":
    main()
